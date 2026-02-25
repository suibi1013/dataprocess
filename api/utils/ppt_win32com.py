from converter.ppt_converter import MsoShapeType
from pptx.dml.color import RGBColor
import re

# Excel边框样式 → PPT COM边框样式映射
# --------------------------
# PPT BorderStyle 枚举值（对应不同边框样式）
border_style_map = {
    "thin": 1,          # 细实线 (ppLineStyleSingle)
    "medium": 2,        # 中粗实线（PPT无直接medium，用宽度区分）
    "thick": 3,         # 粗实线（PPT无直接thick，用宽度区分）
    "double": 4,        # 双实线 (ppLineStyleDouble)
    "dashed": 5,        # 虚线 (ppLineStyleDash)
    "dotted": 6,        # 点线 (ppLineStyleDot)
    "dashDot": 7,       # 点划线 (ppLineStyleDashDot)
    "dashDotDot": 8,    # 双点划线 (ppLineStyleDashDotDot)
    "hair": 1           # 极细（复用thin，通过宽度设为1）
}

# PPT 边框宽度映射（对应Excel的粗细）
border_width_map = {
    "thin": 1,          # 极细/细：1pt
    "medium": 2.25,     # 中粗：2.25pt（PPT默认中等宽度）
    "thick": 3,         # 粗：3pt
    "hair": 0.75,       # 头发丝：0.75pt
    "dashed": 1,        # 虚线默认细
    "dotted": 1,        # 点线默认细
    "dashDot": 1,       # 点划线默认细
    "dashDotDot": 1     # 双点划线默认细
}
style_horizontal={'left':1,'center':2,'right':3}
style_vertical={'top':1,'center':3,'bottom':4}
# 四周边框上1、下2、左3、右4
style_border_map={'top':1,'bottom':2,'left':3,'right':4}
def _hex_to_rgb(hex_str):
    """
    将十六进制RGB字符串（如"FFFF00"、"FF000000"）转换为RGBColor对象
    :param hex_str: 十六进制颜色字符串（6位/8位）
    :return: RGBColor对象
    """
    if hex_str is None:
        return None
    # 去除可能的前缀（如"#"），统一转大写
    hex_str = re.sub(r'^#', '', hex_str).upper()
    
    # 处理8位格式（含Alpha通道，如FF000000），只取后6位
    if len(hex_str) == 8:
        hex_str = hex_str[2:]
    
    # 验证6位格式
    if len(hex_str) != 6 or not re.match(r'^[0-9A-F]{6}$', hex_str):
        raise ValueError(f"无效的RGB格式：{hex_str}，请输入6位十六进制字符串")
    
    # 拆分红、绿、蓝分量并转换为十进制
    r = int(hex_str[0:2], 16)
    g = int(hex_str[2:4], 16)
    b = int(hex_str[4:6], 16)
    
    return RGBColor(r, g, b)
def _rgbcolor_to_int(rgb_color, big_endian=False):
    """
    将RGBColor对象转换为整数类型的RGB值
    :param rgb_color: RGBColor(r, g, b)对象
    :param big_endian: 大端序（True：红高位，符合常规RGB；False：蓝高位，适配部分系统/库）
    :return: 整数类型的RGB值
    """
    # 提取r/g/b分量（0-255）
    if not rgb_color:
        return None
    r, g, b = rgb_color
    
    if big_endian:
        # 大端序（常规RGB）：0xRRGGBB → 例如RGBColor(255,255,0) → 0xFFFF00 = 16776960
        rgb_int = (r << 16) | (g << 8) | b
    else:
        # 小端序（BGR）：0xBBGGRR → 例如RGBColor(255,255,0) → 0x00FFFF = 65535
        rgb_int = (b << 16) | (g << 8) | r
    
    return rgb_int

def _replace_data_msoTable(shape, excel_data: dict):
    """替换表格数据
    Args:
        shape: PPT元素对象（win32com版本）
        excel_data: 要替换的数据和样式
    Returns:
        None
    """
    # 替换表格内容
    if shape.HasTable:
        table = shape.Table
        # 从Excel数据中提取表格数据
        table_data = excel_data['data']
        table_style_map = excel_data['style_map']
        # 替换表格数据
        try:
            # 获取表格的行列数
            num_rows = len(table_data)
            num_cols = len(table_data[0]) if num_rows > 0 else 0
            
            # win32com中表格行列索引从1开始
            for row_idx in range(min(num_rows, table.Rows.Count)):
                row = table.Rows(row_idx + 1)
                for col_idx in range(min(num_cols, row.Cells.Count)):
                    cell = row.Cells(col_idx + 1)
                    if col_idx < len(table_data[row_idx]):
                        cell_value = table_data[row_idx][col_idx]['value']
                        cell_style_id = table_data[row_idx][col_idx]['style_id']
                        text_frame = cell.Shape.TextFrame
                        text_range = cell.Shape.TextFrame.TextRange
                        # === 1. 设置文本内容 ===
                        text_range.Text = cell_value
                        # === 2. 设置样式 ===
                        if cell_style_id in table_style_map:   
                            # === 2.1. 设置字体 ===  
                            style = table_style_map[cell_style_id]
                            if style['font']:
                                font = text_range.Font
                                font.Name = style['font']['name']
                                font.Size = style['font']['size']
                                font.Bold = style['font']['bold']
                                # 确保RGB颜色值是整数类型
                                try:
                                    color_int=_rgbcolor_to_int(_hex_to_rgb(style['font']['color']['rgb']))
                                    if color_int:
                                        font.Color.RGB = _rgbcolor_to_int(_hex_to_rgb(style['font']['color']['rgb']))
                                except Exception as e:
                                    print(f"设置字体颜色失败: {str(e)}")
                                font.Italic = style['font']['italic']
                                font.Underline = style['font']['underline'] if style['font']['underline'] else False
                            # === 2.3. 设置背景色（填充色）===
                            if style['fill']:
                                color_int=_rgbcolor_to_int(_hex_to_rgb(style['fill']['fgColor']['rgb']))
                                if color_int:
                                    fill = cell.Shape.Fill  
                                    fill.Visible = True                              
                                    fill.Solid()                                                 
                                    fill.ForeColor.RGB = color_int

                            # === 2.4. 设置边框 ===
                            if style['border']:
                                for edge in style_border_map:
                                    if edge in style['border']:
                                        # 获取样式和宽度值
                                        border_style = style['border'][edge]['style']
                                        if not border_style:
                                            continue
                                        color_int=_rgbcolor_to_int(_hex_to_rgb(style['border'][edge]['color']['rgb']))
                                        if color_int:                                                    
                                            line_style = border_style_map.get(border_style, 1)
                                            line_width = border_width_map.get(border_style, 1)

                                            border_edge = cell.Borders(style_border_map[edge])
                                            border_edge.Weight = line_width  # 线宽（磅）
                                            border_edge.Style = line_style  # 设置样式（实线/虚线/点线等）
                                            border_edge.ForeColor.RGB = color_int
                                            border_edge.Visible = True
                                        
                            
                            # === 2.5. 设置对齐方式 ===
                            if style['alignment']:
                                # 设置水平对齐方式（例如：居中对齐）
                                horizontal=style['alignment']['horizontal']        
                                if horizontal:                                
                                    text_range.ParagraphFormat.Alignment = style_horizontal[horizontal]
                                # 设置垂直对齐方式（例如：垂直居中对齐）
                                vertical=style['alignment']['vertical']
                                if vertical:
                                    cell.Shape.TextFrame2.VerticalAnchor = style_vertical[vertical]
                        
        except Exception as e:
            raise Exception(f"替换表格数据失败: {str(e)}")
def _replace_data_msoChart(shape, excel_data: dict):
    """替换图表数据
    Args:
        shape: PPT元素对象（win32com版本）
        excel_data: 要替换的数据和样式
    Returns:
        None
    """
    if hasattr(shape, 'Chart'):
        chart = shape.Chart
        # 从Excel数据中提取图表数据
        chart_data = excel_data.get('data', [])
        if not chart_data or len(chart_data) < 2:
            # 至少需要两行数据：第一行为系列名称，第二行开始为类别和数据
            print("图表数据不足，至少需要两行数据：第一行为系列名称，第二行开始为类别和数据")
            return
        
        # 检查并准备数据
        try:
            if not chart_data or len(chart_data) < 2:
                print("图表数据不足，至少需要两行数据：第一行为系列名称，第二行开始为类别和数据")
                return
            
            # 第一行：第一列通常为空或标题，从第二列开始为系列名称
            first_row = chart_data[0]
            if len(first_row) < 2:
                print("图表数据格式不正确，第一行至少需要包含一个系列名称")
                return
            
            # 提取系列名称（从第一行的第二列开始）
            series_names = [item['value'] for item in first_row[1:]] if first_row else []
            if not series_names:
                print("未找到系列名称数据")
                return
            
            # 提取类别名称（从第二行开始，第一列）
            categories = []
            # 准备数据矩阵
            data_matrix = []
            
            # 处理从第二行开始的数据
            for row_idx in range(1, len(chart_data)):
                row = chart_data[row_idx]
                if len(row) < len(first_row):
                    print(f"第{row_idx + 1}行数据长度与第一行不一致，跳过该行")
                    continue
                
                # 第一列是类别名称
                category_name = row[0]['value']
                categories.append(category_name)
                
                # 后续列是数据
                row_data = [item['value'] for item in row[1:]]
                data_matrix.append(row_data)
            
            if not categories:
                print("未找到类别名称数据")
                return
            
            if not data_matrix:
                print("未找到数据部分")
                return
            
            # 确保所有数据行长度一致
            expected_length = len(series_names)
            for row_idx, row_data in enumerate(data_matrix):
                if len(row_data) != expected_length:
                    print(f"第{row_idx + 2}行数据长度与系列名称数量不一致，跳过该行")
                    # 移除对应的类别名称
                    categories.pop(row_idx)
                    # 移除该行数据
                    data_matrix.pop(row_idx)
            
            if not categories or not data_matrix:
                print("没有有效的类别或数据")
                return
            
            # 转换数据格式：按系列组织数据
            series_info_list = []
            for series_idx, series_name in enumerate(series_names):
                # 提取该系列的所有数据点
                series_data = []
                for row_data in data_matrix:
                    if series_idx < len(row_data):
                        series_data.append(row_data[series_idx])
                    else:
                        series_data.append(None)  # 数据缺失
                
                series_info_list.append((series_name, series_data))
            
            print(f"解析图表数据成功，找到 {len(series_info_list)} 个系列，{len(categories)} 个类别")
        except Exception as e:
            print(f"解析图表数据失败: {str(e)}")
            return
        
        # 方法1：更新图表的底层数据源（推荐）
        try:
            # 获取图表的底层数据源（嵌入式Excel工作表）
            data_workbook = chart.ChartData.Workbook
            data_sheet = data_workbook.Worksheets(1)  # 通常是第一个工作表
            
            # 清空所有单元格的内容
            data_sheet.Cells.ClearContents()
            # 同时清空格式，避免保留不必要的单元格格式
            data_sheet.Cells.ClearFormats()
            
            # 写入第一行：第一列留空，后续列是系列名称
            for col_idx, series_name in enumerate(series_names, 2):  # 从第二列开始
                data_sheet.Cells(1, col_idx).Value = series_name
            
            # 写入从第二行开始的数据：第一列是类别名称，后续列是数据
            for row_idx, (category_name, row_data) in enumerate(zip(categories, data_matrix)):
                # 写入类别名称（第一列）
                data_sheet.Cells(row_idx + 2, 1).Value = category_name
                # 写入数据
                for col_idx, value in enumerate(row_data, 2):
                    data_sheet.Cells(row_idx + 2, col_idx).Value = value
                                    
            # 更新图表的数据范围
            try:
                # 计算新的数据范围
                # 第一行：系列名称
                # 第二行开始：类别和数据
                num_categories = len(categories)
                num_series = len(series_names)
                
                if num_categories > 0 and num_series > 0:
                    # 计算数据范围的结束单元格
                    # 第一列是类别，所以数据从第二列开始
                    # 行数：1（系列名称） + num_categories（类别和数据）
                    end_row = 1 + num_categories
                    # 列数：1（类别列） + num_series（数据列）
                    end_col = 1 + num_series
                    
                    # 转换列号为Excel列字母
                    def col_num_to_letter(col_num):
                        letter = ''
                        while col_num > 0:
                            col_num, remainder = divmod(col_num - 1, 26)
                            letter = chr(65 + remainder) + letter
                        return letter
                    
                    end_col_letter = col_num_to_letter(end_col)
                    
                    # 构建数据范围字符串，例如："Sheet1!$A$1:$C$4"
                    data_range = f"Sheet1!$A$1:${end_col_letter}${end_row}"
                    
                    # 更新图表的数据范围,plotBy属性值（1按行、2按列）
                    current_plotby = chart.PlotBy
                    if current_plotby == 0:
                        current_plotby = 1
                    chart.SetSourceData(Source=data_range, PlotBy=current_plotby)
                    print(f"图表数据范围更新成功：{data_range}")
            except Exception as e:
                print(f"更新图表数据范围失败: {str(e)}")
            
            # 刷新图表以显示新数据
            chart.Refresh()
            data_workbook.Close(SaveChanges=True)
            
            print(f"图表底层数据源更新成功，添加了 {len(series_info_list)} 个数据系列")
        except Exception as e:
            print(f"更新图表底层数据源失败: {str(e)}")                        
        
        print(f"图表数据替换成功，添加了 {len(series_info_list)} 个数据系列")
def replace_data_win32com(shape, element_type: str, excel_data: dict):
    """替换表格数据（win32com库版本）
    
    Args:
        shape: PPT元素对象（win32com版本）
        element_type：元素类型
        excel_data: 要替换的数据和样式
    """
    if not excel_data:
        return
    print('excel_data',excel_data)
    if not element_type:
        return
    element_type_int=int(element_type)        
    match element_type_int:
        case MsoShapeType.msoTextBox:
            # 替换文本框内容
            if shape.HasTextFrame:
                text_frame = shape.TextFrame
                text_frame.TextRange.Text = ""  # 清空文本
                text_content = excel_data['data'][0][0]['value']
                text_frame.TextRange.Text = text_content
        case MsoShapeType.msoTable:
            # 替换表格数据
            try:
                _replace_data_msoTable(shape, excel_data)
            except Exception as e:
                raise Exception(f"替换表格数据失败: {str(e)}")
        case MsoShapeType.msoChart:
            # 替换图表数据
            try:
                _replace_data_msoChart(shape, excel_data)
            except Exception as e:
                raise Exception(f"替换图表数据失败: {str(e)}")