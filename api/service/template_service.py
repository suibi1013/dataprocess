#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
模板服务层
负责模板相关的业务逻辑处理
遵循分层架构：服务层调用仓储层，仓储层继承基础仓储
"""

from typing import List, Dict, Any, Optional
from datetime import datetime

from repository.template_info_repository import TemplateInfoRepository
from repository.template_slide_repository import TemplateSlideRepository
from service.base_service import BaseService
from service.result import Result
from config import config
import os

class TemplateService(BaseService):
    """模板服务类"""
    
    def __init__(self, 
                 template_info_repo: TemplateInfoRepository, 
                 template_slide_repo: TemplateSlideRepository):
        """初始化模板服务
        
        Args:
            template_info_repo: 模板信息仓储实例
            template_slide_repo: 模板幻灯片仓储实例
        """
        super().__init__()
        self.template_info_repo = template_info_repo
        self.template_slide_repo = template_slide_repo
        self.upload_folder = config.UPLOAD_FOLDER
    
    async def get_templates(self) -> Result:
        """获取模板列表
        
        Returns:
            Result: 包含模板列表的响应对象
        """
        try:
            # 调用仓储层获取所有模板信息
            template_infos = await self.template_info_repo.find_all()
            
            # 构建响应数据
            templates = []
            for template_info in template_infos:
                template = {
                    'id': template_info.id,
                    'name': template_info.template_name,
                    'file_name': template_info.file_name,
                    'createTime': template_info.created_at,
                    'status': 'ready'
                }
                templates.append(template)
            
            return Result.success(data={
                'templates': templates
            })
        except Exception as e:
            self._log_error(f"获取模板列表失败: {str(e)}")
            return Result.fail(f"获取模板列表失败: {str(e)}")
    
    async def delete_template(self, template_id: str) -> Result:
        """删除指定模板
        
        Args:
            template_id: 模板ID
            
        Returns:
            Result: 删除结果的响应对象
        """
        try:
            # 检查模板是否存在
            template_info = await self.template_info_repo.find_by_id(template_id)
            if not template_info:
                return Result.fail("模板不存在")
            
            # 删除模板幻灯片配置
            await self.template_slide_repo.delete_by_template_id(template_id)
            
            # 删除模板信息
            await self.template_info_repo.delete(template_id)

            # 删除模板目录及文件             
            import os
            import shutil
            template_dir = os.path.join(config.UPLOAD_FOLDER, template_id)
            if os.path.exists(template_dir):
                shutil.rmtree(template_dir, ignore_errors=True)
            
            return Result.success(message="模板删除成功")
        except Exception as e:
            self._log_error(f"删除模板失败: {str(e)}")
            return Result.fail(f"删除模板失败: {str(e)}")
    
    async def check_config_update(self, file_name: str) -> Result:
        """检查配置更新
        
        Args:
            file_name: 文件名
            
        Returns:
            Result: 检查结果的响应对象
        """
        try:
            # 获取所有模板信息
            template_infos = await self.template_info_repo.find_all()
            
            # 查找匹配的模板
            for template_info in template_infos:
                if file_name in template_info.file_name:
                    # 获取最后更新时间戳
                    import time
                    updated_time = datetime.fromisoformat(template_info.updated_at)
                    mtime = time.mktime(updated_time.timetuple())
                    
                    return Result.success(data={
                        'hasUpdate': True,
                        'configFile': template_info.id + '.json',  # 保持与原接口兼容
                        'lastModified': mtime,
                        'message': '找到匹配的配置文件'
                    })
            
            return Result.success(data={
                'hasUpdate': False,
                'message': '未找到匹配的配置文件'
            })
        except Exception as e:
            self._log_error(f"检查配置更新失败: {str(e)}")
            return Result.fail(f"检查配置更新失败: {str(e)}")
    
    async def replace_template_data(self, template_id: str) -> Result:
        """模板数据替换
        
        Args:
            template_id: 模板ID
            
        Returns:
            Result: 数据替换结果的响应对象
        """
        try:
            # 导入必要的模块
            import os
            from utils.excel_helper import ExcelHelper
            from pptx import Presentation
            
            # 获取模板信息
            template_info = await self.template_info_repo.find_by_id(template_id)
            if not template_info:
                return Result.fail("模板不存在")            

            template_file_path = template_info.file_path            
            if not os.path.exists(template_file_path):
                return Result.fail("模板文件不存在")
            
            # 获取模板目录
            template_dir = os.path.dirname(template_file_path)
            
            # 获取模板幻灯片配置
            template_slides = await self.template_slide_repo.find_by_template_id(template_id)
            if not template_slides:
                return Result.fail("模板幻灯片配置不存在")
            
            # 加载PPT文件
            prs = Presentation(template_file_path)
            
            # 遍历每个幻灯片
            for slide_info in template_slides:
                slide_index = slide_info.slide_index
                if slide_index < len(prs.slides):
                    slide = prs.slides[slide_index]
                    elements = slide_info.elements
                    
                    # 遍历每个元素
                    for element in elements:
                        if 'data' in element and 'data_source_config' in element['data']:
                            data_source_config = element['data']['data_source_config']
                            if not data_source_config:
                                continue
                            data_source_path = data_source_config.get('data_source_path')
                            excel_sheet_name = data_source_config.get('excel_sheet_name')
                            excel_cell_range = data_source_config.get('excel_cell_range')
                            
                            if data_source_path and os.path.exists(data_source_path):
                                # 读取Excel数据
                                try:
                                    # 使用ExcelHelper读取数据
                                    excel_data = await ExcelHelper.read_excel_file(
                                        data_source_path, 
                                        excel_sheet_name, 
                                        limit=100
                                    )
                                    
                                    if excel_data['success']:
                                        # 根据元素类型进行数据替换
                                        element_type = element.get('type')
                                        element_id = element.get('id')
                                        
                                        # 遍历PPT中的所有形状，找到对应的元素
                                        for shape in slide.shapes:
                                            # 这里需要根据元素的ID或其他属性找到对应的形状
                                            # 暂时假设元素的id与形状的名称或其他属性对应
                                            # 实际实现需要根据具体的元素标识方式进行调整
                                            if hasattr(shape, 'name') and shape.name == element_id:
                                                # 根据元素类型进行不同的数据替换
                                                if element_type == 'text':
                                                    # 替换文本框内容
                                                    if shape.has_text_frame:
                                                        text_frame = shape.text_frame
                                                        text_frame.clear()
                                                        # 从Excel数据中提取文本内容
                                                        # 这里简化处理，实际需要根据具体的数据结构进行调整
                                                        text_content = self._extract_text_from_excel(excel_data['data'])
                                                        p = text_frame.add_paragraph()
                                                        p.text = text_content
                                                    
                                                elif element_type == 'table':
                                                    # 替换表格内容
                                                    if shape.has_table:
                                                        table = shape.table
                                                        # 从Excel数据中提取表格数据
                                                        table_data = self._extract_table_from_excel(excel_data['data'])
                                                        # 替换表格数据
                                                        self._replace_table_data(table, table_data)
                                                
                                                print(f"成功替换元素 {element_id} 的数据")
                                except Exception as e:
                                    self._log_error(f"读取Excel数据失败: {str(e)}")
            
            # 保存替换后的PPT文件
            # 使用template_info表中的template_name名称
            base_name = template_info.template_name
            ext = '.pptx'  # 确保使用PPTX扩展名
            output_file_name = f"{base_name}{ext}"
            output_file_path = os.path.join(template_dir, output_file_name)
            prs.save(output_file_path)
            
            return Result.success(data={
                'template_id': template_id,
                'output_file_path': output_file_path,
                'message': '数据替换成功'
            })
        except Exception as e:
            self._log_error(f"数据替换失败: {str(e)}")
            return Result.fail(f"数据替换失败: {str(e)}")
    
    def _extract_text_from_excel(self, excel_data: dict) -> str:
        """从Excel数据中提取文本内容
        
        Args:
            excel_data: Excel数据字典
            
        Returns:
            str: 提取的文本内容
        """
        try:
            # 简化处理，实际需要根据具体的数据结构进行调整
            # 这里假设数据结构中包含rows信息
            if 'data' in excel_data:
                data_dict = excel_data['data']
                for key, sheet_data in data_dict.items():
                    if 'rows' in sheet_data and sheet_data['rows']:
                        # 获取第一行第一列的内容作为文本
                        first_row = sheet_data['rows'][0]
                        for col, cell in first_row.items():
                            if 'text' in cell:
                                return cell['text']
            return ""
        except Exception as e:
            self._log_error(f"提取文本内容失败: {str(e)}")
            return ""
    
    def _extract_table_from_excel(self, excel_data: dict) -> list:
        """从Excel数据中提取表格数据
        
        Args:
            excel_data: Excel数据字典
            
        Returns:
            list: 二维表格数据
        """
        try:
            table_data = []
            if 'data' in excel_data:
                data_dict = excel_data['data']
                for key, sheet_data in data_dict.items():
                    if 'rows' in sheet_data:
                        for row in sheet_data['rows']:
                            row_data = []
                            # 按列顺序获取数据
                            columns = sheet_data.get('columns', [])
                            for col in columns:
                                if col in row and 'text' in row[col]:
                                    row_data.append(row[col]['text'])
                                else:
                                    row_data.append("")
                            table_data.append(row_data)
            return table_data
        except Exception as e:
            self._log_error(f"提取表格数据失败: {str(e)}")
            return []
    
    def _replace_table_data(self, table, table_data: list):
        """替换表格数据
        
        Args:
            table: PPT表格对象
            table_data: 二维表格数据
        """
        try:
            # 获取表格的行列数
            num_rows = len(table_data)
            num_cols = len(table_data[0]) if num_rows > 0 else 0
            
            # 遍历表格数据，替换内容
            for row_idx in range(min(num_rows, len(table.rows))):
                row = table.rows[row_idx]
                for col_idx in range(min(num_cols, len(row.cells))):
                    cell = row.cells[col_idx]
                    if col_idx < len(table_data[row_idx]):
                        cell.text = table_data[row_idx][col_idx]
        except Exception as e:
            self._log_error(f"替换表格数据失败: {str(e)}")